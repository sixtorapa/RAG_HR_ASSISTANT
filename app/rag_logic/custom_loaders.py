import logging
import os
import re
from dataclasses import dataclass
from typing import Dict, List, Optional, Tuple

import pdfplumber
from langchain.docstore.document import Document
from langchain.document_loaders.base import BaseLoader
from langchain_community.document_loaders import PyPDFLoader

logging.getLogger("pdfminer").setLevel(logging.ERROR)


# ============================================================
# Loader config
# ============================================================
@dataclass
class LoaderConfig:
    """
    Document loading settings.

    There is no OCR path: neither `pytesseract` nor `tesseract` is installed in
    the deployed images, so any OCR branch would be unreachable and would
    advertise a capability the system does not have. The corpus is text-layer
    PDFs. If scans ever arrive, the dependency is added first.
    """

    pdf_extract_tables: bool = (
        str(os.environ.get("OCR_PDF_EXTRACT_TABLES", "1")).strip().lower() in ("1", "true", "yes")
    )


def _safe_basename(path: str) -> str:
    try:
        return os.path.basename(path)
    except Exception:
        return path


# ============================================================
# PDF loader — never discards a page
# ============================================================
class BetterPDFLoader(BaseLoader):
    """
    Robust PDF loader.
    ✅ Returns one Document per PAGE, never discarding any.
    ✅ Combina: text (fitz, layout-aware) + tablas en Markdown (pdfplumber).
    ✅ Empty pages are marked with is_empty_page=True.
    """

    def __init__(self, file_path: str, loader_cfg: Optional[LoaderConfig] = None):
        self.file_path = file_path
        self.loader_cfg = loader_cfg or LoaderConfig()




    @staticmethod
    def _tables_to_markdown(tables) -> str:
        parts = []
        for table in tables:
            if not table:
                continue
            cleaned_table = [
                [str(cell).replace('\n', ' ').strip() if cell else "" for cell in row]
                for row in table
            ]
            cleaned_table = [row for row in cleaned_table if any(row)]
            if not cleaned_table:
                continue
            md_table = "| " + " | ".join(cleaned_table[0]) + " |\n"
            md_table += "| " + " | ".join(["---"] * len(cleaned_table[0])) + " |\n"
            for row in cleaned_table[1:]:
                md_table += "| " + " | ".join(row) + " |\n"
            parts.append(md_table)
        return "\n\n".join(parts)

    def _load_with_pdfplumber(self) -> List[Document]:
        """
        Fallback for when PyMuPDF (fitz) is unavailable or fails.
        With real text and tables — it must never return [] unless the PDF itself
        is unreadable (corrupt or encrypted).
        """
        fname = _safe_basename(self.file_path)
        try:
            with pdfplumber.open(self.file_path) as pdf:
                page_count = len(pdf.pages)
                if page_count <= 0:
                    return []
                documents: List[Document] = []
                empty_pages = []
                for pidx, page in enumerate(pdf.pages):
                    base_text = (page.extract_text() or "").strip()
                    table_text = ""
                    if self.loader_cfg.pdf_extract_tables:
                        try:
                            tables = page.extract_tables()
                            if tables:
                                table_text = self._tables_to_markdown(tables)
                        except Exception:
                            pass
                    doc_page = self._build_page_document(
                        pidx, page_count, base_text, table_text, fname,
                    )
                    if doc_page.metadata.get("is_empty_page"):
                        empty_pages.append(pidx + 1)
                    documents.append(doc_page)

                print(
                    f"✅ PDF read (pdfplumber fallback) | total_pages={page_count} | "
                    f"pages_with_text={page_count - len(empty_pages)} | empty_pages={len(empty_pages)}"
                )
                if empty_pages:
                    print(f"   📋 Empty pages: {empty_pages}")
                return documents
        except Exception as e:
            print(f"❌ pdfplumber also failed on {fname}: {e}")
            return []

    @staticmethod
    def _strip_table_lines(base_text: str, table_text: str) -> str:
        """
        Remove from the plain text those lines the table already carries in
        markdown.

        The text extractor (fitz/pdfplumber) returns a table as loose words, and
        pdfplumber returns it again as markdown. Without this, the same row
        travels TWICE inside one chunk:

            | 4.5 - 5.0 | Outstanding | Top 10%. Exceptional impact... |
            4.5 - 5.0 Outstanding Top 10%. Exceptional impact...

        Measured before this fix: 11 of the 14 chunks holding a table repeated
        cells, and 56% of the long lines across the whole index were repetitions.
        It is paid in prompt tokens on EVERY query and it distorts BM25 term
        frequencies.

        Deliberately conservative: a line is only removed when its normalised
        form matches a full table row exactly, or a cell of some length. When in
        doubt it is kept — losing real content is far worse than leaving a
        repetition.
        """
        if not (base_text or "").strip() or not (table_text or "").strip():
            return base_text

        def norm(s: str) -> str:
            return re.sub(r"\s+", " ", s or "").strip().lower()

        table_lines = set()
        for line in table_text.splitlines():
            if not line.strip().startswith("|"):
                continue
            cells = [c.strip() for c in line.strip().strip("|").split("|")]
            if all(set(c) <= set("- ") for c in cells):      # row separadora
                continue
            row = norm(" ".join(c for c in cells if c))
            if row:
                table_lines.add(row)
            for c in cells:
                if len(c) >= 12:      # short cells ("25", "Yes") can be legitimate text
                    table_lines.add(norm(c))

        kept = [ln for ln in base_text.splitlines() if norm(ln) not in table_lines]
        # Every removed line leaves a gap. Uncollapsed, dropping three rows in a
        # row leaves six blank lines, paid for in prompt tokens.
        return re.sub(r"\n{3,}", "\n\n", "\n".join(kept)).strip()

    def _build_page_document(self, pidx, page_count, base_text, table_text, fname):
        parts = []

        # 1. Table first when present: it carries the structure
        # The marker stays in Spanish: it is inside the indexed text, so changing
        # it would require a full re-ingest and would invalidate the measured
        # evaluation numbers. It is content, not a comment.
        if (table_text or "").strip():
            parts.append(f"--- TABLE DATA (p. {pidx + 1}) ---\n{table_text.strip()}\n-----------------------------")

        # 2. Base text, with the lines the table just captured removed
        base_text = self._strip_table_lines(base_text, table_text)
        if (base_text or "").strip():
            parts.append(base_text.strip())

        page_out = "\n\n".join(parts).strip()
        is_empty = not page_out or len(page_out.split()) < 3
        if is_empty:
            page_out = f"[Page {pidx + 1} — visual content, no extractable text]"

        return Document(
            page_content=page_out,
            metadata={
                "source": self.file_path, "file_type": "pdf",
                "filename": fname, "source_file": fname, "relative_path": fname,
                "page": pidx, "page_number": pidx + 1, "page_count": page_count,
                "has_table": bool((table_text or "").strip()),
                "is_empty_page": is_empty, "text_chars": len(page_out),
            },
        )

    def load(self) -> List[Document]:
        fname = _safe_basename(self.file_path)
        print(f"📄 Processing PDF: {fname}")
        documents: List[Document] = []

        # === NIVEL 1: PyMuPDF ===
        try:
            import fitz
            with fitz.open(self.file_path) as doc:
                page_count = doc.page_count
                if page_count <= 0:
                    return []
                per_page_text, per_page_visual = [], []
                for page in doc:
                    # sort=True to respect columns and visual layout
                    t = (page.get_text("text", sort=True) or "").strip()
                    
                    has_images = bool(page.get_images(full=True))
                    try:
                        has_drawings = bool(page.get_drawings())
                    except Exception:
                        has_drawings = False
                    per_page_text.append(t)
                    per_page_visual.append(bool(has_images or has_drawings))

                total_chars = sum(len(t) for t in per_page_text)

                # Table extraction to Markdown
                table_map = {}
                if self.loader_cfg.pdf_extract_tables:
                    try:
                        with pdfplumber.open(self.file_path) as plumber_pdf:
                            for pidx in range(min(page_count, len(plumber_pdf.pages))):
                                try:
                                    tables = plumber_pdf.pages[pidx].extract_tables()
                                    if tables:
                                        md = self._tables_to_markdown(tables)
                                        if md:
                                            table_map[pidx] = md
                                except Exception:
                                    continue
                    except Exception as e:
                        print(f"   ⚠️ Table extraction failed: {e}")

                empty_pages = []
                for pidx in range(page_count):
                    doc_page = self._build_page_document(
                        pidx, page_count, per_page_text[pidx],
                        table_map.get(pidx, ""), fname,
                    )
                    if doc_page.metadata.get("is_empty_page"):
                        empty_pages.append(pidx + 1)
                    documents.append(doc_page)

                print(
                    f"✅ PDF read (PyMuPDF Enhanced) | total_pages={page_count} | "
                    f"pages_with_text={page_count - len(empty_pages)} | pages_with_tables={len(table_map)} | "
                    f"empty_pages={len(empty_pages)} | total_chars={total_chars}"
                )
                if empty_pages:
                    print(f"   📋 Empty pages: {empty_pages}")
                return documents
        except Exception as e:
            print(f"⚠️ PyMuPDF failed on {fname}: {e} — falling back to pdfplumber")

        # === LEVEL 2: pdfplumber — never return [] merely because fitz is missing ===
        return self._load_with_pdfplumber()


# ============================================================
# PowerPoint loader — never discards a slide
# ============================================================
class BetterPowerPointLoader(BaseLoader):
    def __init__(self, file_path: str, loader_cfg: Optional[LoaderConfig] = None):
        self.file_path = file_path
        self.loader_cfg = loader_cfg or LoaderConfig()

    def _extract_slide_text(self, slide):
        parts = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    t = (shape.text or "").strip()
                    if t:
                        parts.append(t)
            except Exception:
                continue
        return "\n".join(parts).strip()


    def load(self):
        fname = _safe_basename(self.file_path)
        print(f"📊 Processing PPT: {fname}")
        try:
            from pptx import Presentation
        except Exception as e:
            print(f"❌ python-pptx missing: {e}")
            return []
        try:
            prs = Presentation(self.file_path)
        except Exception as e:
            print(f"❌ Could not open PPT: {e}")
            return []

        out_docs = []
        slide_count = len(prs.slides)
        for s_i, slide in enumerate(prs.slides, start=1):
            slide_text = self._extract_slide_text(slide)
            page_out = slide_text.strip()
            is_empty = not page_out or len(page_out.split()) < 3
            if is_empty:
                page_out = f"[Slide {s_i} — visual content, no extractable text]"
            out_docs.append(Document(page_content=page_out, metadata={
                "source": self.file_path, "file_type": "ppt",
                "filename": fname, "source_file": fname, "relative_path": fname,
                "slide": s_i, "slide_number": s_i, "slide_count": slide_count,
                "is_empty_slide": is_empty, "text_chars": len(page_out),
            }))
        print("✅ PPT read")
        return out_docs